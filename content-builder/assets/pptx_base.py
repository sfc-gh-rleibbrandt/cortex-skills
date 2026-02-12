#!/usr/bin/env python3
"""
Content Builder - PPTX Base Template
Creates professional PowerPoint presentations with Snowflake-style branding.

This template includes all the patterns and helper functions developed for
customer briefings, with consistent styling across all slides.

Usage:
    1. Copy this file to your project
    2. Customize the SLIDES section at the bottom
    3. Run: python3 create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# =============================================================================
# SNOWFLAKE BRAND COLORS
# =============================================================================
class Colors:
    """Snowflake brand colors for consistent styling."""
    SF_BLUE = RGBColor(41, 181, 232)        # #29B5E8 - Accent, links
    SF_DARK_BLUE = RGBColor(17, 86, 127)    # #11567F - Titles, headers
    SF_NAVY = RGBColor(13, 44, 84)          # #0D2C54 - Table headers
    SF_LIGHT_BG = RGBColor(244, 250, 255)   # #F4FAFF - Backgrounds
    SF_GRAY = RGBColor(110, 118, 129)       # #6E7681 - Subtitles
    TABLE_HEADER = RGBColor(41, 128, 185)   # Blue for table headers
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(26, 26, 26)
    LIGHT_GRAY = RGBColor(100, 100, 100)
    RED = RGBColor(180, 50, 50)             # For "before" examples
    GREEN = RGBColor(50, 150, 50)           # For "after" examples
    HIGHLIGHT_GREEN = RGBColor(39, 174, 96) # For filled values

# =============================================================================
# PRESENTATION SETUP
# =============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Get layouts
blank_layout = prs.slide_layouts[6]  # Blank - used for all custom slides

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_copyright_footer(slide, light=False):
    """Add copyright footer to slide."""
    footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(8), Inches(0.3))
    tf = footer_box.text_frame
    tf.paragraphs[0].text = "© 2026 Snowflake Inc. All Rights Reserved | Confidential"
    tf.paragraphs[0].font.size = Pt(9)
    if light:
        tf.paragraphs[0].font.color.rgb = RGBColor(160, 196, 232)
    else:
        tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)


def add_doc_link(slide, url, display_text=None):
    """Add documentation link to bottom of slide."""
    if display_text is None:
        display_text = url.replace("https://", "")
    link_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
    tf = link_box.text_frame
    tf.paragraphs[0].text = f"📚 {display_text}"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = Colors.LIGHT_GRAY


def add_title_slide(title, subtitle="", event_name="", date_presenter=""):
    """
    Add a title slide with dark blue background.
    Used for presentation title and section dividers.
    """
    slide = prs.slides.add_slide(blank_layout)
    
    # Dark blue background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.SF_DARK_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = Colors.WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = RGBColor(232, 244, 252)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Event name
    if event_name:
        event_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6))
        tf = event_box.text_frame
        tf.paragraphs[0].text = event_name
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(232, 244, 252)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Date/Presenter
    if date_presenter:
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5))
        tf = date_box.text_frame
        tf.paragraphs[0].text = date_presenter
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RGBColor(160, 196, 232)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_copyright_footer(slide, light=True)
    return slide


def add_section_slide(title, subtitle=""):
    """Add a section divider slide with dark blue background."""
    slide = prs.slides.add_slide(blank_layout)
    
    # Dark blue background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = Colors.SF_DARK_BLUE
    bg.line.fill.background()
    
    # Title - centered
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = Colors.WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.8))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(232, 244, 252)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_copyright_footer(slide, light=True)
    return slide


def add_content_slide(title, subtitle=""):
    """
    Add a content slide with fixed header position.
    Returns the slide for adding additional content.
    """
    slide = prs.slides.add_slide(blank_layout)
    
    # Title at fixed position
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = Colors.SF_DARK_BLUE
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.4))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = Colors.SF_GRAY
    
    add_copyright_footer(slide)
    return slide


def add_text_content(slide, content, left=0.5, top=1.4, width=12, height=5):
    """Add text content to a slide."""
    content_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = content
    tf.paragraphs[0].font.size = Pt(14)
    return content_box


def add_table(slide, data, left=0.5, top=1.5, width=12, row_height=0.4):
    """
    Add a styled table to a slide.
    First row is treated as header (blue background, white text).
    """
    rows = len(data)
    cols = len(data[0])
    
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), 
                                          Inches(width), Inches(row_height * rows))
    table = table_shape.table
    
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            
            if i == 0:  # Header row
                para.font.bold = True
                para.font.color.rgb = Colors.WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = Colors.TABLE_HEADER
            else:
                para.font.color.rgb = Colors.BLACK
    
    return table


def add_code_block(slide, code, left=0.5, top=1.5, width=12, height=3, font_size=11):
    """Add a code block with monospace font."""
    code_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = code_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = code
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.name = "Courier New"
    return code_box


def add_before_after_slide(title, problem, before_code, after_code, 
                           comparison_data=None, customer_quote=None, doc_url=None):
    """
    Add a before/after comparison slide.
    This is a key pattern for showing value of new features.
    """
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    # Problem statement
    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    tf = problem_box.text_frame
    tf.paragraphs[0].text = f"Problem: {problem}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    
    # BEFORE label
    before_label = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(6), Inches(0.3))
    tf = before_label.text_frame
    tf.paragraphs[0].text = "❌ BEFORE: Complex, error-prone, slow"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = Colors.RED
    
    # BEFORE code
    before_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(6.2), Inches(2.4))
    tf = before_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = before_code
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.name = "Courier New"
    
    # AFTER label
    after_label = slide.shapes.add_textbox(Inches(6.7), Inches(1.1), Inches(6), Inches(0.3))
    tf = after_label.text_frame
    tf.paragraphs[0].text = "✅ AFTER: Clean, fast, correct"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = Colors.GREEN
    
    # AFTER code
    after_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(6), Inches(1.5))
    tf = after_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = after_code
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.name = "Courier New"
    
    # Comparison table (if provided)
    if comparison_data:
        table = slide.shapes.add_table(len(comparison_data), len(comparison_data[0]), 
                                       Inches(6.7), Inches(3.0), Inches(6), Inches(1.5)).table
        for i, row_data in enumerate(comparison_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                if i == 0:
                    para.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Colors.TABLE_HEADER
                    para.font.color.rgb = Colors.WHITE
    
    # Customer quote (if provided)
    if customer_quote:
        quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.5))
        tf = quote_box.text_frame
        tf.paragraphs[0].text = f"🏆 {customer_quote}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
    
    # Doc link
    if doc_url:
        add_doc_link(slide, doc_url)
    
    add_copyright_footer(slide)
    return slide


def add_value_prop_slide(title, subtitle, features_table, struggles, values, customer_quote=None, doc_url=None):
    """
    Add a value proposition slide with:
    - Feature table
    - Common struggles we solve
    - Business value delivered
    - Optional customer quote
    """
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    # Subtitle (value prop)
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.4))
    tf = sub_box.text_frame
    tf.paragraphs[0].text = subtitle
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = Colors.SF_BLUE
    
    # Features table
    add_table(slide, features_table, top=1.2, width=12.2, row_height=0.35)
    
    # Common struggles label
    struggle_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(6), Inches(0.3))
    tf = struggle_label.text_frame
    tf.paragraphs[0].text = "Common Struggles We Solve:"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    
    # Struggles list
    struggle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(5.5), Inches(1.5))
    tf = struggle_box.text_frame
    for i, struggle in enumerate(struggles):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"❌ {struggle}"
        p.font.size = Pt(12)
    
    # Value label
    value_label = slide.shapes.add_textbox(Inches(7), Inches(3.6), Inches(6), Inches(0.3))
    tf = value_label.text_frame
    tf.paragraphs[0].text = "Business Value Delivered:"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    
    # Values list
    value_box = slide.shapes.add_textbox(Inches(7), Inches(3.9), Inches(5.5), Inches(1.5))
    tf = value_box.text_frame
    for i, value in enumerate(values):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✅ {value}"
        p.font.size = Pt(12)
    
    # Customer quote
    if customer_quote:
        quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12), Inches(0.5))
        tf = quote_box.text_frame
        tf.paragraphs[0].text = f"🏆 {customer_quote}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
    
    if doc_url:
        add_doc_link(slide, doc_url)
    
    add_copyright_footer(slide)
    return slide


def add_roadmap_slide(title, pain_point, current_behavior, roadmap_table, planned_syntax=None, doc_url=None):
    """
    Add a roadmap slide for upcoming features.
    """
    slide = prs.slides.add_slide(blank_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    
    # Pain point quote
    pain_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.8))
    tf = pain_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = f'"{pain_point}"'
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True
    
    # Current behavior
    current_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(0.4))
    tf = current_box.text_frame
    tf.paragraphs[0].text = f"Current: {current_behavior}"
    tf.paragraphs[0].font.size = Pt(14)
    
    # Roadmap table
    add_table(slide, roadmap_table, top=2.2, width=6, row_height=0.4)
    
    # Planned syntax
    if planned_syntax:
        syntax_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(6), Inches(0.3))
        tf = syntax_label.text_frame
        tf.paragraphs[0].text = "Planned Syntax:"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        
        add_code_block(slide, planned_syntax, top=4.4, width=8, height=1.5, font_size=12)
    
    if doc_url:
        add_doc_link(slide, doc_url)
    
    add_copyright_footer(slide)
    return slide


# =============================================================================
# PRESENTATION CONTENT - CUSTOMIZE BELOW
# =============================================================================

# Example: Title slide
add_title_slide(
    title="Presentation Title",
    subtitle="Subtitle Here",
    event_name="Customer/Event Name",
    date_presenter="January 2026 | Presenter Name"
)

# Example: Section slide
add_section_slide("Part 1: Introduction", "Overview and Context")

# Example: Content slide with table
slide = add_content_slide("Feature Overview", "Key capabilities")
add_table(slide, [
    ["Feature", "Status", "Description"],
    ["Feature A", "✅ GA", "Available today"],
    ["Feature B", "🟡 Preview", "Coming soon"],
    ["Feature C", "🔜 Roadmap", "Q1 FY27"],
], top=1.4)
add_doc_link(slide, "https://docs.snowflake.com")

# Example: Before/After comparison
add_before_after_slide(
    title="New Feature - Before vs After ✅ GA",
    problem="Complex queries required multiple CTEs and window functions",
    before_code="""WITH complex_cte AS (
  SELECT *, ROW_NUMBER() OVER (
    PARTITION BY id ORDER BY ts DESC
  ) AS rn
  FROM table1 t1
  LEFT JOIN (
    SELECT *, LEAD(ts) OVER (
      PARTITION BY id ORDER BY ts
    ) AS next_ts
    FROM table2
  ) t2 ON t1.id = t2.id
    AND t1.ts >= t2.ts
    AND t1.ts < t2.next_ts
)
SELECT * FROM complex_cte WHERE rn = 1;""",
    after_code="""SELECT * FROM table1 t1
NEW_JOIN table2 t2
  MATCH_CONDITION(t1.ts >= t2.ts)
  ON t1.id = t2.id;""",
    comparison_data=[
        ["Metric", "Before", "After"],
        ["Lines of SQL", "15+", "5"],
        ["Readability", "Complex", "Intuitive"],
        ["Performance", "Slow", "Optimized"],
    ],
    customer_quote='"99% performance improvement over previous solutions" — Customer Name',
    doc_url="https://docs.snowflake.com/path/to/docs"
)

# Example: Value proposition slide
add_value_prop_slide(
    title="Platform Capabilities",
    subtitle="No separate specialized database needed → Lower TCO, unified governance",
    features_table=[
        ["Function", "What It Solves", "Without It..."],
        ["Feature A", "Problem A", "Manual workaround"],
        ["Feature B", "Problem B", "Complex alternative"],
        ["Feature C", "Problem C", "Tedious approach"],
    ],
    struggles=[
        "Pain point 1",
        "Pain point 2",
        "Pain point 3",
        "Pain point 4",
    ],
    values=[
        "Benefit 1",
        "Benefit 2",
        "Benefit 3",
        "Benefit 4",
    ],
    customer_quote='"Customer success quote" — Customer Name',
    doc_url="https://docs.snowflake.com"
)

# Example: Roadmap slide
add_roadmap_slide(
    title="Upcoming Feature",
    pain_point="Customer pain point that this feature solves",
    current_behavior="How things work today (with limitations)",
    roadmap_table=[
        ["Milestone", "Timeline"],
        ["Design", "Q4 FY26 (Complete)"],
        ["Preview", "Q1 FY27"],
        ["GA", "Q2 FY27"],
    ],
    planned_syntax="""CREATE NEW_FEATURE feature_name
    SETTING = 'value'
AS SELECT ...;""",
    doc_url="https://docs.snowflake.com"
)

# Example: Thank you slide
add_title_slide(
    title="Thank You!",
    subtitle="Questions?",
    event_name="",
    date_presenter="Contact: email@snowflake.com | Next: Topic @ Time"
)

# =============================================================================
# SAVE
# =============================================================================
output_path = "presentation_output.pptx"
prs.save(output_path)
print(f"✅ Created: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
